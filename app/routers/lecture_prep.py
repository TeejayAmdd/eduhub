"""
Lecture Prep with Cortex — generates a downloadable PPTX from lecture material.
"""

import io
import os
import uuid
import logging
from typing import Optional

import fitz  # PyMuPDF
from fastapi import APIRouter, Depends, HTTPException, UploadFile, File, Form
from fastapi.responses import Response
from sqlalchemy.orm import Session, joinedload

import app.models as models
from app.database import get_db
from app.auth import get_current_user

log = logging.getLogger(__name__)
router = APIRouter(prefix="/api/lecture-prep", tags=["Lecture Prep"])

MAX_MATERIAL_CHARS = 60_000

# In-memory store: presentation_id -> pptx bytes (TTL handled by session lifecycle)
_store: dict[str, bytes] = {}
# Parallel store: presentation_id -> structured slide data from Claude (for preview)
_store_data: dict[str, dict] = {}


# ── pptx imports ──────────────────────────────────────────────────────────────
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.chart.data import ChartData
    from pptx.enum.chart import XL_CHART_TYPE
    _PPTX_OK = True
except ImportError:
    _PPTX_OK = False

# ── Theme & font definitions ───────────────────────────────────────────────────
THEMES: dict[str, dict] = {
    "dark": {
        "title_bg":   (0x0F, 0x17, 0x2A), "slide_bg":  (0x0F, 0x17, 0x2A),
        "accent":     (0x4F, 0x46, 0xE5), "title_txt": (0xFF, 0xFF, 0xFF),
        "hdr_txt":    (0xFF, 0xFF, 0xFF), "body_txt":  (0xE2, 0xE8, 0xF0),
        "muted":      (0x94, 0xA3, 0xB8), "dot":       (0x4F, 0x46, 0xE5),
    },
    "corporate": {
        "title_bg":   (0x00, 0x32, 0x74), "slide_bg":  (0xFF, 0xFF, 0xFF),
        "accent":     (0x00, 0x71, 0xBC), "title_txt": (0xFF, 0xFF, 0xFF),
        "hdr_txt":    (0xFF, 0xFF, 0xFF), "body_txt":  (0x1A, 0x1A, 0x1A),
        "muted":      (0x5A, 0x6B, 0x8A), "dot":       (0x00, 0x71, 0xBC),
    },
    "minimal": {
        "title_bg":   (0x11, 0x11, 0x11), "slide_bg":  (0xFA, 0xFA, 0xF9),
        "accent":     (0x11, 0x11, 0x11), "title_txt": (0xFF, 0xFF, 0xFF),
        "hdr_txt":    (0xFF, 0xFF, 0xFF), "body_txt":  (0x22, 0x22, 0x22),
        "muted":      (0x77, 0x77, 0x77), "dot":       (0x11, 0x11, 0x11),
    },
    "forest": {
        "title_bg":   (0x0D, 0x2B, 0x18), "slide_bg":  (0xF0, 0xFD, 0xF4),
        "accent":     (0x16, 0xA3, 0x4A), "title_txt": (0xFF, 0xFF, 0xFF),
        "hdr_txt":    (0xFF, 0xFF, 0xFF), "body_txt":  (0x14, 0x53, 0x2D),
        "muted":      (0x16, 0xA3, 0x4A), "dot":       (0x16, 0xA3, 0x4A),
    },
    "sunset": {
        "title_bg":   (0x18, 0x08, 0x02), "slide_bg":  (0xFF, 0xF7, 0xED),
        "accent":     (0xEA, 0x58, 0x0C), "title_txt": (0xFF, 0xFF, 0xFF),
        "hdr_txt":    (0xFF, 0xFF, 0xFF), "body_txt":  (0x43, 0x18, 0x07),
        "muted":      (0x97, 0x4E, 0x12), "dot":       (0xEA, 0x58, 0x0C),
    },
}

FONTS: dict[str, str] = {
    "calibri":   "Calibri",
    "arial":     "Arial",
    "georgia":   "Georgia",
    "trebuchet": "Trebuchet MS",
}

# ── Low-level drawing helpers ─────────────────────────────────────────────────
def _rgb(t): return RGBColor(*t)


def _rect(slide, l, t, w, h, fill=None):
    shp = slide.shapes.add_shape(1, l, t, w, h)
    if fill:
        shp.fill.solid()
        shp.fill.fore_color.rgb = _rgb(fill)
    else:
        shp.fill.background()
    shp.line.fill.background()
    return shp


def _txt(slide, l, t, w, h, text, size, bold=False, color=(0xFF,0xFF,0xFF), align=None, font="Calibri"):
    from pptx.enum.text import PP_ALIGN
    align = align or PP_ALIGN.LEFT
    txb = slide.shapes.add_textbox(l, t, w, h)
    txb.word_wrap = True
    tf = txb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = _rgb(color)
    run.font.name = font
    return txb


def _add_logo(slide, logo_bytes: bytes | None, W, H):
    if not logo_bytes:
        return
    try:
        slide.shapes.add_picture(io.BytesIO(logo_bytes), W - Inches(1.1), H - Inches(0.45), Inches(1.0), Inches(0.35))
    except Exception:
        pass


def _slide_bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = _rgb(color)


def _content_header(slide, title, theme, font, W):
    _rect(slide, Inches(0), Inches(0), W, Inches(1.2), theme["accent"])
    _txt(slide, Inches(0.4), Inches(0.15), Inches(12.5), Inches(0.9),
         title, 28, bold=True, color=theme["hdr_txt"], font=font)


def _slide_number(slide, num, theme, font, W, H):
    _txt(slide, W - Inches(0.7), H - Inches(0.38), Inches(0.6), Inches(0.3),
         str(num), 11, color=theme["muted"], font=font)


# ── Per-type slide builders ────────────────────────────────────────────────────
def _bullets_slide(prs, blank, info, num, theme, font, logo_bytes, W, H):
    s = prs.slides.add_slide(blank)
    _slide_bg(s, theme["slide_bg"])
    _content_header(s, info.get("title", ""), theme, font, W)
    bullets = info.get("bullets", [])[:6]
    y0, bh = Inches(1.45), Inches(0.73)
    for i, b in enumerate(bullets):
        y = y0 + i * bh
        _rect(s, Inches(0.52), y + Inches(0.28), Inches(0.11), Inches(0.11), theme["dot"])
        _txt(s, Inches(0.78), y, Inches(12.2), bh, b, 18, color=theme["body_txt"], font=font)
    if info.get("notes"):
        s.notes_slide.notes_text_frame.text = info["notes"]
    _slide_number(s, num, theme, font, W, H)
    _add_logo(s, logo_bytes, W, H)


def _chart_slide(prs, blank, info, num, theme, font, logo_bytes, W, H):
    s = prs.slides.add_slide(blank)
    _slide_bg(s, theme["slide_bg"])
    _content_header(s, info.get("title", ""), theme, font, W)
    cd_raw = info.get("chart_data", {})
    cats = cd_raw.get("categories", [])
    series_list = cd_raw.get("series", [])
    ct_map = {
        "bar":    XL_CHART_TYPE.BAR_CLUSTERED,
        "column": XL_CHART_TYPE.COLUMN_CLUSTERED,
        "line":   XL_CHART_TYPE.LINE_MARKERS,
        "pie":    XL_CHART_TYPE.PIE,
    }
    ct = ct_map.get(info.get("chart_type", "column"), XL_CHART_TYPE.COLUMN_CLUSTERED)
    if cats and series_list:
        cd = ChartData()
        cd.categories = cats
        for sr in series_list:
            cd.add_series(sr.get("name", ""), sr.get("values", []))
        gf = s.shapes.add_chart(ct, Inches(0.5), Inches(1.4), Inches(12.3), Inches(5.6), cd)
        chart = gf.chart
        chart.has_legend = len(series_list) > 1
    else:
        _txt(s, Inches(0.5), Inches(2), Inches(12), Inches(1),
             "Chart data unavailable", 18, color=theme["muted"], font=font)
    if info.get("notes"):
        s.notes_slide.notes_text_frame.text = info["notes"]
    _slide_number(s, num, theme, font, W, H)
    _add_logo(s, logo_bytes, W, H)


def _two_column_slide(prs, blank, info, num, theme, font, logo_bytes, W, H):
    s = prs.slides.add_slide(blank)
    _slide_bg(s, theme["slide_bg"])
    _content_header(s, info.get("title", ""), theme, font, W)
    col_w = Inches(6.0)
    # Left column
    lh = info.get("left_heading", "")
    if lh:
        _rect(s, Inches(0.4), Inches(1.35), Inches(6.0), Inches(0.04), theme["accent"])
        _txt(s, Inches(0.4), Inches(1.4), col_w, Inches(0.5), lh, 16, bold=True, color=theme["accent"], font=font)
    for i, b in enumerate(info.get("left_bullets", [])[:5]):
        y = Inches(1.95) + i * Inches(0.7)
        _rect(s, Inches(0.48), y + Inches(0.22), Inches(0.09), Inches(0.09), theme["dot"])
        _txt(s, Inches(0.65), y, Inches(5.7), Inches(0.65), b, 16, color=theme["body_txt"], font=font)
    # Divider
    _rect(s, Inches(6.65), Inches(1.3), Inches(0.03), Inches(5.8), theme["muted"])
    # Right column
    rh = info.get("right_heading", "")
    if rh:
        _rect(s, Inches(6.8), Inches(1.35), Inches(6.0), Inches(0.04), theme["accent"])
        _txt(s, Inches(6.8), Inches(1.4), col_w, Inches(0.5), rh, 16, bold=True, color=theme["accent"], font=font)
    for i, b in enumerate(info.get("right_bullets", [])[:5]):
        y = Inches(1.95) + i * Inches(0.7)
        _rect(s, Inches(6.88), y + Inches(0.22), Inches(0.09), Inches(0.09), theme["dot"])
        _txt(s, Inches(7.05), y, Inches(5.7), Inches(0.65), b, 16, color=theme["body_txt"], font=font)
    if info.get("notes"):
        s.notes_slide.notes_text_frame.text = info["notes"]
    _slide_number(s, num, theme, font, W, H)
    _add_logo(s, logo_bytes, W, H)


def _quote_slide(prs, blank, info, num, theme, font, logo_bytes, W, H):
    s = prs.slides.add_slide(blank)
    _slide_bg(s, theme["title_bg"])
    # Accent bar left
    _rect(s, Inches(0), Inches(0), Inches(0.5), H, theme["accent"])
    # Opening quote mark
    _txt(s, Inches(0.9), Inches(0.8), Inches(2), Inches(1.5), "“", 80, bold=True,
         color=theme["accent"], font=font)
    # Quote text
    qt = info.get("quote_text", "")
    _txt(s, Inches(1.2), Inches(1.8), Inches(11.0), Inches(4.0), qt, 26, bold=False,
         color=theme["title_txt"], font=font)
    # Source
    src = info.get("quote_source", "")
    if src:
        _txt(s, Inches(1.2), Inches(5.9), Inches(11.0), Inches(0.5),
             f"— {src}", 16, color=theme["muted"], font=font)
    if info.get("notes"):
        s.notes_slide.notes_text_frame.text = info["notes"]
    _slide_number(s, num, theme, font, W, H)
    _add_logo(s, logo_bytes, W, H)


def _divider_slide(prs, blank, info, num, theme, font, logo_bytes, W, H):
    s = prs.slides.add_slide(blank)
    _slide_bg(s, theme["accent"])
    from pptx.enum.text import PP_ALIGN
    _txt(s, Inches(1.5), Inches(2.8), Inches(10.3), Inches(1.8),
         info.get("title", ""), 44, bold=True, color=(0xFF, 0xFF, 0xFF),
         align=PP_ALIGN.CENTER, font=font)
    _rect(s, Inches(4.5), Inches(4.7), Inches(4.3), Inches(0.06), (0xFF, 0xFF, 0xFF))
    _slide_number(s, num, theme, font, W, H)
    _add_logo(s, logo_bytes, W, H)


def _image_placeholder_slide(prs, blank, info, num, theme, font, logo_bytes, W, H):
    s = prs.slides.add_slide(blank)
    _slide_bg(s, theme["slide_bg"])
    _content_header(s, info.get("title", ""), theme, font, W)
    # Dashed placeholder box
    box = s.shapes.add_shape(1, Inches(0.7), Inches(1.5), Inches(11.9), Inches(5.4))
    box.fill.background()
    box.line.color.rgb = _rgb(theme["muted"])
    box.line.width = Pt(1.5)
    # Description text inside box
    desc = info.get("image_description", "Insert image or diagram here")
    from pptx.enum.text import PP_ALIGN
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = f"[ {desc} ]"
    run.font.size = Pt(18)
    run.font.color.rgb = _rgb(theme["muted"])
    run.font.name = font
    run.font.italic = True
    if info.get("notes"):
        s.notes_slide.notes_text_frame.text = info["notes"]
    _slide_number(s, num, theme, font, W, H)
    _add_logo(s, logo_bytes, W, H)


# ── Master builder ─────────────────────────────────────────────────────────────
def _build_pptx(data: dict, theme_name: str = "dark", font_name: str = "calibri",
                logo_bytes: bytes | None = None) -> bytes:
    if not _PPTX_OK:
        raise HTTPException(500, "python-pptx not installed on this server")

    theme = THEMES.get(theme_name, THEMES["dark"])
    font  = FONTS.get(font_name, "Calibri")

    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    W, H  = Inches(13.33), Inches(7.5)

    title_text = data.get("title", "Lecture Presentation")
    subject    = data.get("subject", "")
    slides     = data.get("slides", [])

    # ── Title slide ──────────────────────────────────────────────────────────
    ts = prs.slides.add_slide(blank)
    _slide_bg(ts, theme["title_bg"])
    _rect(ts, Inches(0), Inches(0), Inches(0.55), H, theme["accent"])
    _rect(ts, Inches(0.55), Inches(5.2), Inches(12.78), Inches(0.05), theme["accent"])
    from pptx.enum.text import PP_ALIGN
    _txt(ts, Inches(1), Inches(2.2), Inches(11.5), Inches(1.7),
         title_text, 44, bold=True, color=theme["title_txt"], font=font)
    if subject:
        _txt(ts, Inches(1), Inches(4.0), Inches(11.5), Inches(0.7),
             subject, 22, color=theme["muted"], font=font)
    _txt(ts, Inches(1), Inches(6.7), Inches(6), Inches(0.45),
         "Generated by Cortex AI", 11, color=theme["muted"], font=font)
    _add_logo(ts, logo_bytes, W, H)

    # ── Content slides ────────────────────────────────────────────────────────
    builders = {
        "bullets":           _bullets_slide,
        "chart":             _chart_slide,
        "two_column":        _two_column_slide,
        "quote":             _quote_slide,
        "divider":           _divider_slide,
        "image_placeholder": _image_placeholder_slide,
    }
    for idx, info in enumerate(slides):
        slide_type = info.get("slide_type", "bullets")
        builder = builders.get(slide_type, _bullets_slide)
        builder(prs, blank, info, idx + 2, theme, font, logo_bytes, W, H)

    out = io.BytesIO()
    prs.save(out)
    return out.getvalue()




# ── Claude slide generation ───────────────────────────────────────────────────

_SLIDES_TOOL = {
    "name": "create_presentation",
    "description": "Create a structured lecture presentation from the provided material.",
    "input_schema": {
        "type": "object",
        "properties": {
            "title": {"type": "string", "description": "Concise presentation title"},
            "subject": {"type": "string", "description": "One-line course/topic description for the subtitle"},
            "slides": {
                "type": "array",
                "description": "Content slides — title slide is auto-generated, do not include it here.",
                "items": {
                    "type": "object",
                    "properties": {
                        "title": {"type": "string", "description": "Slide heading, max 8 words"},
                        "slide_type": {
                            "type": "string",
                            "enum": ["bullets", "chart", "two_column", "quote", "divider", "image_placeholder"],
                            "description": (
                                "bullets: standard bullet list | "
                                "chart: data visualisation (bar/column/line/pie) | "
                                "two_column: side-by-side comparison | "
                                "quote: highlighted definition or key quote | "
                                "divider: section break (full-colour) | "
                                "image_placeholder: slot for a diagram or image the lecturer will insert"
                            ),
                        },
                        "bullets": {
                            "type": "array", "items": {"type": "string"},
                            "description": "3-5 bullet points, max 20 words each. Required when slide_type is 'bullets'.",
                        },
                        "chart_type": {
                            "type": "string", "enum": ["bar", "column", "line", "pie"],
                            "description": "Required when slide_type is 'chart'. bar=horizontal, column=vertical, line=trend, pie=proportions.",
                        },
                        "chart_data": {
                            "type": "object",
                            "description": "Required when slide_type is 'chart'. Use real numbers extracted from the material.",
                            "properties": {
                                "categories": {"type": "array", "items": {"type": "string"}},
                                "series": {
                                    "type": "array",
                                    "items": {
                                        "type": "object",
                                        "properties": {
                                            "name":   {"type": "string"},
                                            "values": {"type": "array", "items": {"type": "number"}},
                                        },
                                        "required": ["name", "values"],
                                    },
                                },
                            },
                            "required": ["categories", "series"],
                        },
                        "left_heading":  {"type": "string",  "description": "Left column heading for 'two_column' slides"},
                        "right_heading": {"type": "string",  "description": "Right column heading for 'two_column' slides"},
                        "left_bullets":  {"type": "array", "items": {"type": "string"}, "description": "2-4 bullets for left column"},
                        "right_bullets": {"type": "array", "items": {"type": "string"}, "description": "2-4 bullets for right column"},
                        "quote_text":    {"type": "string",  "description": "The quote or definition, max 50 words. Required for 'quote' slides."},
                        "quote_source":  {"type": "string",  "description": "Attribution or textbook reference for 'quote' slides"},
                        "image_description": {"type": "string", "description": "Detailed description of the diagram/image for 'image_placeholder' slides"},
                        "notes": {"type": "string", "description": "Lecturer speaker notes, 1-3 sentences"},
                    },
                    "required": ["title", "slide_type"],
                },
            },
        },
        "required": ["title", "subject", "slides"],
    },
}


def _call_claude(material_text: str, prefs: dict) -> dict:
    api_key = os.environ.get("ANTHROPIC_API_KEY", "")
    if not api_key:
        raise HTTPException(500, "AI service not configured")

    import anthropic
    client = anthropic.Anthropic(api_key=api_key)

    num_slides       = prefs.get("num_slides", "auto")
    level            = prefs.get("level", "undergraduate")
    focus            = prefs.get("focus", "key concepts and main points")
    tone             = prefs.get("tone", "academic")
    notes            = prefs.get("include_notes", True)
    include_charts   = prefs.get("include_charts", False)
    include_diagrams = prefs.get("include_diagrams", False)
    enhanced_layouts = prefs.get("enhanced_layouts", False)

    slide_target = {
        "10-15": "10 to 15", "15-20": "15 to 20",
        "20-30": "20 to 30", "auto": "an appropriate number (usually 12-18)",
    }.get(num_slides, "an appropriate number")

    # Build conditional slide-type guidance
    chart_rule = (
        "- Use 'chart' slides (1-3 max) wherever the material contains clear numerical or comparative data "
        "(statistics, percentages, trends). Extract real numbers from the text and pick the best chart type."
        if include_charts else
        "- Do NOT use 'chart' slide_type."
    )
    diagram_rule = (
        "- Use 'image_placeholder' slides (1-2 max) where a diagram, figure, or process flow would "
        "significantly aid understanding. Write a detailed image_description so the lecturer knows exactly what to insert."
        if include_diagrams else
        "- Do NOT use 'image_placeholder' slide_type."
    )
    layout_rules = (
        "- Use 'two_column' slides (1-3 max) for comparisons, pros/cons, before/after, or two contrasting concepts.\n"
        "- Use 'quote' slides (1-2 max) for important definitions, key theorems, or memorable statements.\n"
        "- Use 'divider' slides (1-2 max) only when the presentation has clear major sections (15+ slides)."
        if enhanced_layouts else
        "- Do NOT use 'two_column', 'quote', or 'divider' slide types — use 'bullets' for all content slides."
    )

    system = f"""You are an expert educational content designer creating professional lecture slides.

Extract and organise the lecture material into a clear, well-structured presentation.

PREFERENCES:
- Slides: {slide_target} content slides (not counting the title slide)
- Student level: {level}
- Emphasis: {focus if focus else "key concepts and main points"}
- Tone: {tone}
- Speaker notes: {"include thoughtful speaker notes for each slide" if notes else "omit speaker notes"}

SLIDE TYPE RULES:
{chart_rule}
{diagram_rule}
{layout_rules}
- Default slide_type is 'bullets' — use it for all standard content slides.
- Start with an Introduction/Overview slide (bullets).
- End with a Summary or Key Takeaways slide (bullets).
- Slide titles: concise, max 8 words.
- Bullet points: 3-5 per slide, max 20 words each, no sub-bullets.
- Do NOT include the title slide in the slides array — it is generated automatically."""

    response = client.messages.create(
        model="claude-haiku-4-5-20251001",
        max_tokens=6000,
        system=system,
        messages=[{"role": "user", "content": f"<lecture_material>\n{material_text}\n</lecture_material>"}],
        tools=[_SLIDES_TOOL],
        tool_choice={"type": "tool", "name": "create_presentation"},
    )

    for block in response.content:
        if block.type == "tool_use" and block.name == "create_presentation":
            return block.input

    raise HTTPException(500, "AI failed to generate slide content. Please try again.")


# ── Text extraction ───────────────────────────────────────────────────────────

def _extract_text(data: bytes, mime: str) -> str:
    if "pdf" in mime or mime == "application/octet-stream":
        try:
            doc = fitz.open(stream=data, filetype="pdf")
            return "\n".join(p.get_text() for p in doc)[:MAX_MATERIAL_CHARS]
        except Exception:
            pass
    # Plain text / docx fallback
    try:
        return data.decode("utf-8", errors="ignore")[:MAX_MATERIAL_CHARS]
    except Exception:
        return ""


# ── Endpoints ─────────────────────────────────────────────────────────────────

@router.post("/generate")
async def generate_slides(
    file: UploadFile = File(...),
    title: str = Form(""),
    num_slides: str = Form("auto"),
    level: str = Form("undergraduate"),
    focus: str = Form(""),
    tone: str = Form("academic"),
    include_notes: bool = Form(True),
    theme: str = Form("dark"),
    font: str = Form("calibri"),
    include_charts: bool = Form(False),
    include_diagrams: bool = Form(False),
    enhanced_layouts: bool = Form(False),
    logo: Optional[UploadFile] = File(None),
    db: Session = Depends(get_db),
    current_user: models.User = Depends(get_current_user),
):
    if current_user.role != models.RoleEnum.lecturer:
        raise HTTPException(403, "Only lecturers can generate presentations")

    raw = await file.read()
    if not raw:
        raise HTTPException(400, "Uploaded file is empty")

    material_text = _extract_text(raw, file.content_type or "")
    if not material_text.strip():
        raise HTTPException(400, "Could not extract text from the file. Please use a PDF or plain-text document.")

    logo_bytes: Optional[bytes] = None
    if logo and logo.filename:
        try:
            logo_bytes = await logo.read() or None
        except Exception:
            pass

    prefs = {
        "num_slides": num_slides,
        "level": level,
        "focus": focus,
        "tone": tone,
        "include_notes": include_notes,
        "include_charts": include_charts,
        "include_diagrams": include_diagrams,
        "enhanced_layouts": enhanced_layouts,
    }

    try:
        slide_data = _call_claude(material_text, prefs)
    except HTTPException:
        raise
    except Exception as exc:
        log.error("Claude error in lecture-prep: %s", exc, exc_info=True)
        raise HTTPException(500, f"AI error: {exc}")

    if title.strip():
        slide_data["title"] = title.strip()

    try:
        pptx_bytes = _build_pptx(slide_data, theme_name=theme, font_name=font, logo_bytes=logo_bytes)
    except Exception as exc:
        log.error("PPTX build error: %s", exc, exc_info=True)
        raise HTTPException(500, f"Slide build error: {exc}")

    presentation_id = str(uuid.uuid4())
    _store[presentation_id] = pptx_bytes
    _store_data[presentation_id] = slide_data

    content_slides = slide_data.get("slides", [])

    def _preview_body(s: dict) -> list[str]:
        st = s.get("slide_type", "bullets")
        if st == "bullets":    return s.get("bullets", [])
        if st == "two_column": return (s.get("left_bullets") or []) + (s.get("right_bullets") or [])
        if st == "quote":      return [s.get("quote_text", "")]
        if st == "chart":
            cats = s.get("chart_data", {}).get("categories", [])
            return [f"Chart: {s.get('chart_type','column')} — {', '.join(str(c) for c in cats[:4])}"]
        if st == "image_placeholder": return [f"[Image: {s.get('image_description','')}]"]
        if st == "divider":    return []
        return s.get("bullets", [])

    preview_slides = [
        {"slide": 1, "title": slide_data.get("title", ""), "body": [slide_data.get("subject", "")], "notes": [],
         "slide_type": "bullets"},
        *[
            {
                "slide": i + 2,
                "title": s.get("title", ""),
                "slide_type": s.get("slide_type", "bullets"),
                "body": _preview_body(s),
                "notes": [s["notes"]] if s.get("notes") else [],
            }
            for i, s in enumerate(content_slides)
        ],
    ]

    return {
        "presentation_id": presentation_id,
        "title": slide_data.get("title", ""),
        "subject": slide_data.get("subject", ""),
        "slide_count": len(content_slides) + 1,
        "outline": [s.get("title", "") for s in content_slides],
        "slides": preview_slides,
    }


@router.get("/download/{presentation_id}")
def download_presentation(
    presentation_id: str,
    current_user: models.User = Depends(get_current_user),
):
    pptx = _store.get(presentation_id)
    if not pptx:
        raise HTTPException(404, "Presentation not found or expired. Please generate again.")

    return Response(
        content=pptx,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers={"Content-Disposition": f'attachment; filename="cortex-lecture.pptx"'},
    )


@router.post("/save/{class_id}")
def save_to_course(
    class_id: int,
    payload: dict,
    db: Session = Depends(get_db),
    current_user: models.User = Depends(get_current_user),
):
    if current_user.role != models.RoleEnum.lecturer:
        raise HTTPException(403, "Only lecturers can save materials")

    presentation_id = payload.get("presentation_id", "")
    custom_title    = (payload.get("title") or "").strip()

    pptx = _store.get(presentation_id)
    if not pptx:
        raise HTTPException(404, "Presentation not found or expired. Please generate again.")

    cls = db.query(models.Class).filter(
        models.Class.id == class_id,
        models.Class.lecturer_id == current_user.id,
    ).first()
    if not cls:
        raise HTTPException(404, "Course not found or you don't own it")

    mat_title = custom_title or f"Lecture Slides — {cls.name}"
    filename  = f"cortex-slides-{uuid.uuid4().hex[:8]}.pptx"

    material = models.CourseMaterial(
        class_id         = class_id,
        uploaded_by      = current_user.id,
        title            = mat_title,
        description      = "Generated by Cortex AI Lecture Prep",
        original_filename= filename,
        stored_filename  = filename,
        file_type        = "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        file_size        = len(pptx),
        file_data        = pptx,
    )
    db.add(material)
    db.commit()
    db.refresh(material)

    # Also write to personal history so it appears without needing a course
    sd = _store_data.get(presentation_id)
    history_entry = models.LecturePrepHistory(
        user_id=current_user.id,
        title=mat_title,
        original_filename=filename,
        file_data=pptx,
        class_name=cls.name,
        course_code=cls.course_code,
        slide_data=sd,
    )
    db.add(history_entry)
    db.commit()

    return {"message": f"Saved to {cls.name}", "material_id": material.id}


@router.post("/save-history")
def save_to_history(
    payload: dict,
    db: Session = Depends(get_db),
    current_user: models.User = Depends(get_current_user),
):
    """Save a generated presentation to personal history without attaching it to a course."""
    if current_user.role != models.RoleEnum.lecturer:
        raise HTTPException(403, "Only lecturers can save presentations")

    presentation_id = payload.get("presentation_id", "")
    custom_title = (payload.get("title") or "").strip()

    pptx = _store.get(presentation_id)
    if not pptx:
        raise HTTPException(404, "Presentation not found or expired. Please generate again.")

    title = custom_title or "Lecture Slides"
    filename = f"cortex-slides-{uuid.uuid4().hex[:8]}.pptx"

    entry = models.LecturePrepHistory(
        user_id=current_user.id,
        title=title,
        original_filename=filename,
        file_data=pptx,
        slide_data=_store_data.get(presentation_id),
    )
    db.add(entry)
    db.commit()
    db.refresh(entry)
    return {"message": "Saved to history", "history_id": entry.id}


@router.get("/history")
def get_history(
    db: Session = Depends(get_db),
    current_user: models.User = Depends(get_current_user),
):
    """Returns all lecture prep presentations saved by the user, most recent first."""
    rows = (
        db.query(models.LecturePrepHistory)
        .filter(models.LecturePrepHistory.user_id == current_user.id)
        .order_by(models.LecturePrepHistory.saved_at.desc())
        .limit(50)
        .all()
    )
    return [
        {
            "history_id": r.id,
            "title": r.title,
            "class_name": r.class_name,
            "course_code": r.course_code,
            "saved_at": r.saved_at.isoformat() if r.saved_at else None,
        }
        for r in rows
    ]


@router.get("/material/{material_id}/download")
def redownload_material(
    material_id: int,
    db: Session = Depends(get_db),
    current_user: models.User = Depends(get_current_user),
):
    """Re-download a previously saved lecture prep PPTX."""
    mat = db.query(models.CourseMaterial).filter(
        models.CourseMaterial.id == material_id,
        models.CourseMaterial.uploaded_by == current_user.id,
    ).first()
    if not mat or not mat.file_data:
        raise HTTPException(404, "Material not found")
    return Response(
        content=mat.file_data,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers={"Content-Disposition": f'attachment; filename="{mat.original_filename}"'},
    )


@router.get("/history/{history_id}/slides")
def get_history_slides(
    history_id: int,
    db: Session = Depends(get_db),
    current_user: models.User = Depends(get_current_user),
):
    """Return structured slide data for in-browser preview of a saved presentation."""
    entry = db.query(models.LecturePrepHistory).filter(
        models.LecturePrepHistory.id == history_id,
        models.LecturePrepHistory.user_id == current_user.id,
    ).first()
    if not entry:
        raise HTTPException(404, "History item not found")
    if not entry.slide_data:
        raise HTTPException(404, "No preview data available — re-generate the presentation to get a preview")

    sd = entry.slide_data
    content_slides = sd.get("slides", [])
    return [
        {"slide": 1, "title": sd.get("title", ""), "body": [sd.get("subject", "")], "notes": []},
        *[
            {
                "slide": i + 2,
                "title": s.get("title", ""),
                "body": s.get("bullets", []),
                "notes": [s["notes"]] if s.get("notes") else [],
            }
            for i, s in enumerate(content_slides)
        ],
    ]


@router.get("/history/{history_id}/download")
def redownload_history(
    history_id: int,
    db: Session = Depends(get_db),
    current_user: models.User = Depends(get_current_user),
):
    """Re-download a PPTX from personal lecture prep history."""
    entry = db.query(models.LecturePrepHistory).filter(
        models.LecturePrepHistory.id == history_id,
        models.LecturePrepHistory.user_id == current_user.id,
    ).first()
    if not entry or not entry.file_data:
        raise HTTPException(404, "History item not found")
    return Response(
        content=entry.file_data,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers={"Content-Disposition": f'attachment; filename="{entry.original_filename}"'},
    )
